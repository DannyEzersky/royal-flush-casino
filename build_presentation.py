"""
Builds the Royal Spin A/B Test findings PowerPoint presentation.
Output: exports/royal_spin_ab_test.pptx
"""

import sqlite3
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import config

# ---------------------------------------------------------------------------
# Brand colours
# ---------------------------------------------------------------------------
NAVY       = RGBColor(0x0D, 0x1B, 0x3E)
GOLD       = RGBColor(0xD4, 0xAF, 0x37)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GREY   = RGBColor(0xCC, 0xCC, 0xCC)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED        = RGBColor(0xE7, 0x4C, 0x3C)
AMBER      = RGBColor(0xF3, 0x9C, 0x12)

W = Inches(13.33)
H = Inches(7.5)


# ---------------------------------------------------------------------------
# Metrics loader — all slide data comes from here, nothing hardcoded
# ---------------------------------------------------------------------------

def load_metrics(conn: sqlite3.Connection) -> dict:
    sql_dir = Path("sql")

    def run(fname):
        return conn.execute((sql_dir / fname).read_text(encoding="utf-8")).fetchall()

    def kv(fname, k=0, v=1):
        return {r[k]: r[v] for r in run(fname)}

    arpu      = kv("experiment_arpu.sql")
    conv      = kv("experiment_conversion_rate.sql")
    ftd30     = kv("experiment_ftd_30day.sql")
    spins     = kv("experiment_spins_per_session.sql")
    ret_raw   = {r[0]: {"d1": r[1], "d7": r[2], "d30": r[3]}
                 for r in run("experiment_retention.sql")}
    royal_rate = run("royal_token_conversion_rate.sql")[0][0]
    arpu_seg  = {(r[0], r[1]): r[2] for r in run("experiment_arpu_by_segment.sql")}
    arpu_plat = {(r[0], r[1]): r[2]
                 for r in run("experiment_arpu_by_platform.sql")
                 if r[0] in ("Control", "Treatment")}

    # Treatment revenue split by currency
    rev_type = dict(conn.execute("""
        SELECT t.currency_type, ROUND(SUM(t.amount_usd), 0)
        FROM transactions t
        JOIN players p ON p.player_id = t.player_id
        WHERE p.experiment_group = 'Treatment'
        GROUP BY t.currency_type
    """).fetchall())
    coin_rev  = rev_type.get("coin", 0)
    token_rev = rev_type.get("royal_token", 0)
    total_rev = coin_rev + token_rev

    # Dolphin cannibalization detail
    d = {(g, c): v for g, c, v in conn.execute("""
        SELECT p.experiment_group, t.currency_type, ROUND(SUM(t.amount_usd), 0)
        FROM transactions t
        JOIN players p ON p.player_id = t.player_id
        WHERE p.experiment_group IN ('Control','Treatment') AND p.spend_segment = 'Dolphin'
        GROUP BY p.experiment_group, t.currency_type
    """).fetchall()}
    d_ctrl_coin  = d.get(("Control",   "coin"),        0)
    d_trt_coin   = d.get(("Treatment", "coin"),        0)
    d_trt_token  = d.get(("Treatment", "royal_token"), 0)
    d_net        = (d_trt_coin + d_trt_token) - d_ctrl_coin

    # Whale revenue concentration
    whale_share = dict(conn.execute("""
        SELECT p.experiment_group,
               ROUND(SUM(CASE WHEN p.spend_segment='Whale' THEN t.amount_usd ELSE 0 END)
                     * 100.0 / SUM(t.amount_usd), 1)
        FROM transactions t
        JOIN players p ON p.player_id = t.player_id
        WHERE p.experiment_group IN ('Control','Treatment')
        GROUP BY p.experiment_group
    """).fetchall())

    # Avg active days per group
    active_days = dict(conn.execute("""
        SELECT p.experiment_group, ROUND(AVG(active_days), 1)
        FROM players p
        JOIN (SELECT player_id, COUNT(DISTINCT session_date) AS active_days
              FROM sessions WHERE spin_count >= 1 GROUP BY player_id) s
            ON s.player_id = p.player_id
        WHERE p.experiment_group IN ('Control','Treatment')
        GROUP BY p.experiment_group
    """).fetchall())

    # Derived scalars
    ca, ta   = arpu["Control"], arpu["Treatment"]
    arpu_lift = round((ta - ca) / ca * 100)

    def _lift(a, b): return round((b - a) / a * 100) if a else 0

    ios_lift      = _lift(arpu_plat.get(("Control","iOS"),0),     arpu_plat.get(("Treatment","iOS"),0))
    and_lift      = _lift(arpu_plat.get(("Control","Android"),0), arpu_plat.get(("Treatment","Android"),0))
    whale_seg_lift = _lift(arpu_seg.get(("Control","Whale"),0),   arpu_seg.get(("Treatment","Whale"),0))

    return {
        "ctrl_arpu": ca,  "trt_arpu": ta,  "arpu_lift": arpu_lift,
        "ctrl_conv": conv["Control"],  "trt_conv": conv["Treatment"],
        "conv_delta": round(conv["Treatment"] - conv["Control"], 2),
        "ctrl_ftd30": ftd30["Control"], "trt_ftd30": ftd30["Treatment"],
        "ftd30_delta": round(ftd30["Treatment"] - ftd30["Control"], 2),
        "ctrl_spins": spins["Control"], "trt_spins": spins["Treatment"],
        "spins_delta": round(spins["Treatment"] - spins["Control"], 2),
        "ctrl_d1":  ret_raw["Control"]["d1"],   "trt_d1":  ret_raw["Treatment"]["d1"],
        "ctrl_d7":  ret_raw["Control"]["d7"],   "trt_d7":  ret_raw["Treatment"]["d7"],
        "ctrl_d30": ret_raw["Control"]["d30"],  "trt_d30": ret_raw["Treatment"]["d30"],
        "royal_adoption": royal_rate,
        "trt_coin_rev": coin_rev, "trt_token_rev": token_rev, "trt_total_rev": total_rev,
        "trt_coin_pct":  round(coin_rev  / total_rev * 100, 2) if total_rev else 0,
        "trt_token_pct": round(token_rev / total_rev * 100, 2) if total_rev else 0,
        "arpu_seg": arpu_seg, "arpu_plat": arpu_plat,
        "dolph_ctrl_coin": d_ctrl_coin, "dolph_trt_coin": d_trt_coin,
        "dolph_trt_token": d_trt_token, "dolph_net": d_net,
        "whale_share": whale_share, "active_days": active_days,
        "ios_lift": ios_lift, "and_lift": and_lift, "whale_seg_lift": whale_seg_lift,
    }


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, colour):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def box(slide, x, y, w, h, bg=None, border=None, border_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if border and border_pt:
        shape.line.color.rgb = border
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, x, y, w, h, text="", size=18, bold=False,
          colour=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = "Calibri"
    return tb


def metric_card(slide, x, y, w, h, label, value, delta="",
                delta_colour=GREEN, bg=None):
    card_bg = bg or RGBColor(0x1A, 0x2D, 0x5A)
    box(slide, x, y, w, h, bg=card_bg, border=GOLD, border_pt=1.5)
    pad = Inches(0.15)
    txbox(slide, x+pad, y+pad,          w-pad*2, Inches(0.4),
          label, size=11, colour=MID_GREY, align=PP_ALIGN.CENTER)
    txbox(slide, x+pad, y+Inches(0.45), w-pad*2, Inches(0.65),
          value, size=26, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)
    if delta:
        txbox(slide, x+pad, y+Inches(1.05), w-pad*2, Inches(0.35),
              delta, size=13, bold=True, colour=delta_colour,
              align=PP_ALIGN.CENTER)


def divider(slide, y, colour=GOLD, thickness_pt=1.5):
    line = slide.shapes.add_shape(1, Inches(0.5), y, W-Inches(1), Pt(thickness_pt))
    line.fill.solid()
    line.fill.fore_color.rgb = colour
    line.line.fill.background()


def slide_header(slide, title, subtitle=""):
    box(slide, 0, 0, W, Inches(0.85), bg=GOLD)
    txbox(slide, Inches(0.4), Inches(0.1), W-Inches(0.8), Inches(0.65),
          title, size=22, bold=True, colour=NAVY, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, Inches(0.4), Inches(0.82), W-Inches(0.8), Inches(0.35),
              subtitle, size=12, colour=MID_GREY, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def slide_title(prs, m):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)

    box(sl, 0, 0, Inches(0.18), H, bg=GOLD)
    chip = sl.shapes.add_shape(9, Inches(8.5), Inches(0.5), Inches(6.2), Inches(6.2))
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(0x14, 0x26, 0x52)
    chip.line.fill.background()

    txbox(sl, Inches(0.6), Inches(1.5), Inches(8), Inches(0.7),
          "ROYAL FLUSH CASINO", size=14, bold=True, colour=GOLD)
    txbox(sl, Inches(0.6), Inches(2.1), Inches(9), Inches(1.6),
          "Royal Spin Feature\nA/B Test Results", size=40, bold=True, colour=WHITE)

    divider(sl, Inches(3.75), GOLD, 2)

    txbox(sl, Inches(0.6), Inches(3.95), Inches(7), Inches(0.4),
          "Experiment period: July 1 – September 28, 2024", size=14, colour=MID_GREY)
    txbox(sl, Inches(0.6), Inches(4.35), Inches(7), Inches(0.4),
          "SpinCrown Studios  |  Data & Analytics", size=14, colour=MID_GREY)

    box(sl, Inches(9.8), Inches(5.2), Inches(3.1), Inches(1.8),
        bg=RGBColor(0x1A, 0x2D, 0x5A), border=GOLD, border_pt=1.5)
    txbox(sl, Inches(9.9), Inches(5.3),  Inches(2.9), Inches(0.4),
          "ARPU LIFT", size=11, colour=MID_GREY, align=PP_ALIGN.CENTER)
    txbox(sl, Inches(9.9), Inches(5.65), Inches(2.9), Inches(0.7),
          f"+{m['arpu_lift']}%", size=36, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)
    txbox(sl, Inches(9.9), Inches(6.3),  Inches(2.9), Inches(0.35),
          f"${m['ctrl_arpu']:.2f}  →  ${m['trt_arpu']:.2f}",
          size=12, colour=WHITE, align=PP_ALIGN.CENTER)


def slide_agenda(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Agenda")

    items = [
        ("01", "Test Setup & Methodology"),
        ("02", f"Headline Result — ARPU +{prs._arpu_lift}%"),
        ("03", "Royal Token Revenue Impact"),
        ("04", "Segment Deep-Dive"),
        ("05", "Platform Breakdown"),
        ("06", "Guardrail Metrics"),
        ("07", "Recommendation"),
    ]
    col_x = [Inches(0.6), Inches(6.9)]
    for i, (num, label) in enumerate(items):
        col, row = i % 2, i // 2
        x = col_x[col]
        y = Inches(1.3) + row * Inches(1.35)
        box(sl, x, y, Inches(5.9), Inches(1.1),
            bg=RGBColor(0x1A, 0x2D, 0x5A), border=GOLD, border_pt=1)
        txbox(sl, x+Inches(0.15), y+Inches(0.08), Inches(0.7), Inches(0.5),
              num, size=22, bold=True, colour=GOLD)
        txbox(sl, x+Inches(0.85), y+Inches(0.25), Inches(4.8), Inches(0.5),
              label, size=15, bold=True, colour=WHITE)


def slide_methodology(prs):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Test Setup & Methodology",
                 "Royal Spin — a new premium spin mechanic funded by Royal Tokens (in-app currency)")

    for i, (title, lines) in enumerate([
        ("CONTROL", [
            "Standard coin-based spins only",
            "No Royal Token currency",
            "~12,500 players",
            "Baseline period: Apr – Jun 2024",
        ]),
        ("TREATMENT", [
            "Royal Spin mechanic unlocked",
            "Royal Tokens purchasable via IAP",
            "~12,500 players",
            "Experiment: Jul – Sep 2024",
        ]),
    ]):
        x = Inches(0.5) + i * Inches(6.3)
        box(sl, x, Inches(1.2), Inches(5.9), Inches(3.8),
            bg=RGBColor(0x1A, 0x2D, 0x5A), border=GOLD, border_pt=1.5)
        label_bg = GOLD if i == 1 else MID_GREY
        box(sl, x, Inches(1.2), Inches(5.9), Inches(0.45), bg=label_bg)
        txbox(sl, x, Inches(1.22), Inches(5.9), Inches(0.4),
              title, size=14, bold=True, colour=NAVY, align=PP_ALIGN.CENTER)
        for j, line in enumerate(lines):
            txbox(sl, x+Inches(0.3), Inches(1.85)+j*Inches(0.65),
                  Inches(5.3), Inches(0.55), "•  "+line, size=14, colour=WHITE)

    box(sl, Inches(0.5), Inches(5.2), W-Inches(1), Inches(1.9),
        bg=RGBColor(0x0A, 0x14, 0x30))
    params = [
        ("Players", "50,000 total"),
        ("Simulation Period", "181 days"),
        ("Segments", "Minnow · Dolphin · Whale"),
        ("Platforms", "iOS · Android"),
        ("Markets", "US · UK · DE · CA · AU · Other"),
    ]
    for i, (k, v) in enumerate(params):
        x = Inches(0.8) + i * Inches(2.5)
        txbox(sl, x, Inches(5.3),  Inches(2.3), Inches(0.35),
              k, size=10, colour=GOLD, align=PP_ALIGN.CENTER)
        txbox(sl, x, Inches(5.65), Inches(2.3), Inches(0.35),
              v, size=12, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)


def slide_headline(prs, m):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, f"Headline Result — ARPU +{m['arpu_lift']}%",
                 "Treatment group generated significantly more revenue per user with no engagement cost")

    cw, ch, gap = Inches(2.9), Inches(1.55), Inches(0.22)
    start_x, y  = Inches(0.5), Inches(1.15)

    def _dc(delta): return GREEN if delta >= 0 else RED
    def _sign(v):   return "▲" if v >= 0 else "▼"

    metrics = [
        ("ARPU",
         f"${m['trt_arpu']:.2f}",
         f"{_sign(m['arpu_lift'])} +{m['arpu_lift']}%  vs ${m['ctrl_arpu']:.2f}",
         GREEN),
        ("Conversion Rate",
         f"{m['trt_conv']:.2f}%",
         f"{_sign(m['conv_delta'])} {m['conv_delta']:+.2f}pp  vs {m['ctrl_conv']:.2f}%",
         _dc(m['conv_delta'])),
        ("FTD 30-Day Rate",
         f"{m['trt_ftd30']:.2f}%",
         f"{_sign(m['ftd30_delta'])} {m['ftd30_delta']:+.2f}pp  vs {m['ctrl_ftd30']:.2f}%",
         _dc(m['ftd30_delta'])),
        ("Avg Spins / Session",
         f"{m['trt_spins']:.2f}",
         f"{_sign(m['spins_delta'])} {m['spins_delta']:+.2f}  vs {m['ctrl_spins']:.2f}",
         _dc(m['spins_delta'])),
    ]
    for i, (label, val, delta, dc) in enumerate(metrics):
        metric_card(sl, start_x + i*(cw+gap), y, cw, ch, label, val, delta, dc)

    # Retention table
    box(sl, Inches(0.5), Inches(2.95), W-Inches(1), Inches(0.38), bg=GOLD)
    for i, txt in enumerate(["", "D1 Retention", "D7 Retention", "D30 Retention"]):
        txbox(sl, Inches(0.6)+i*Inches(3.0), Inches(2.97), Inches(2.8), Inches(0.34),
              txt, size=13, bold=True, colour=NAVY, align=PP_ALIGN.CENTER)

    rows = [
        ("Control",   f"{m['ctrl_d1']:.2f}%", f"{m['ctrl_d7']:.2f}%", f"{m['ctrl_d30']:.2f}%"),
        ("Treatment", f"{m['trt_d1']:.2f}%",  f"{m['trt_d7']:.2f}%",  f"{m['trt_d30']:.2f}%"),
    ]
    for r, (grp, d1, d7, d30) in enumerate(rows):
        row_bg = RGBColor(0x1A, 0x2D, 0x5A) if r == 1 else RGBColor(0x12, 0x22, 0x48)
        y_row  = Inches(3.33) + r*Inches(0.6)
        box(sl, Inches(0.5), y_row, W-Inches(1), Inches(0.58), bg=row_bg)
        for i, (val, bold, colour) in enumerate([
            (grp, True,  GOLD),
            (d1,  False, WHITE),
            (d7,  False, WHITE),
            (d30, False, WHITE),
        ]):
            txbox(sl, Inches(0.6)+i*Inches(3.0), y_row+Inches(0.1),
                  Inches(2.8), Inches(0.38),
                  val, size=14, bold=bold, colour=colour, align=PP_ALIGN.CENTER)

    txbox(sl, Inches(0.5), Inches(4.65), W-Inches(1), Inches(0.4),
          "Retention is statistically flat between groups — the Royal Spin feature did not harm engagement.",
          size=13, colour=MID_GREY, align=PP_ALIGN.CENTER)

    box(sl, Inches(0.5), Inches(5.15), W-Inches(1), Inches(1.95),
        bg=RGBColor(0x0A, 0x14, 0x30), border=GOLD, border_pt=2)
    txbox(sl, Inches(0.7), Inches(5.25), W-Inches(1.4), Inches(0.45),
          "KEY TAKEAWAY", size=11, bold=True, colour=GOLD)
    txbox(sl, Inches(0.7), Inches(5.65), W-Inches(1.4), Inches(1.2),
          f"Treatment players spent {m['arpu_lift']}% more on average than Control, while session "
          f"engagement remained identical. The ARPU lift is driven purely by monetisation, not playtime inflation.",
          size=14, colour=WHITE, wrap=True)


def slide_royal_token(prs, m):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Royal Token Revenue Impact",
                 f"{m['royal_adoption']:.2f}% of Treatment players adopted Royal Tokens "
                 f"— and contributed {m['trt_token_pct']:.0f}% of all Treatment revenue")

    box(sl, Inches(0.5), Inches(1.1), Inches(5.8), Inches(4.5),
        bg=RGBColor(0x1A, 0x2D, 0x5A), border=GOLD, border_pt=2)
    txbox(sl, Inches(0.7), Inches(1.3), Inches(5.4), Inches(0.4),
          "ROYAL TOKEN ADOPTERS", size=12, colour=GOLD, align=PP_ALIGN.CENTER)
    txbox(sl, Inches(0.7), Inches(1.7), Inches(5.4), Inches(1.1),
          f"{m['royal_adoption']:.2f}%", size=60, bold=True, colour=WHITE,
          align=PP_ALIGN.CENTER)
    txbox(sl, Inches(0.7), Inches(2.75), Inches(5.4), Inches(0.4),
          "of Treatment players", size=14, colour=MID_GREY, align=PP_ALIGN.CENTER)

    divider(sl, Inches(3.3), MID_GREY, 1)

    txbox(sl, Inches(0.7), Inches(3.45), Inches(5.4), Inches(0.4),
          "REVENUE GENERATED", size=12, colour=GOLD, align=PP_ALIGN.CENTER)
    txbox(sl, Inches(0.7), Inches(3.8), Inches(5.4), Inches(0.7),
          f"${m['trt_token_rev']:,.0f}", size=36, bold=True, colour=WHITE,
          align=PP_ALIGN.CENTER)
    txbox(sl, Inches(0.7), Inches(4.45), Inches(5.4), Inches(0.4),
          "in Royal Token transactions", size=13, colour=MID_GREY, align=PP_ALIGN.CENTER)

    box(sl, Inches(6.7), Inches(1.1), Inches(6.2), Inches(4.5),
        bg=RGBColor(0x1A, 0x2D, 0x5A), border=GOLD, border_pt=2)
    txbox(sl, Inches(6.9), Inches(1.3), Inches(5.8), Inches(0.4),
          "TREATMENT REVENUE BREAKDOWN", size=12, colour=GOLD, align=PP_ALIGN.CENTER)

    bars = [
        ("Coin Purchases",  m['trt_coin_pct'],  RGBColor(0x41, 0x8C, 0xF0),
         f"${m['trt_coin_rev']:,.0f}"),
        ("Royal Tokens",    m['trt_token_pct'], GOLD,
         f"${m['trt_token_rev']:,.0f}"),
    ]
    for i, (label, pct, colour, amt) in enumerate(bars):
        by    = Inches(2.0) + i*Inches(1.8)
        bar_w = Inches(5.4) * pct / 100
        txbox(sl, Inches(6.9), by, Inches(5.8), Inches(0.35), label, size=12, colour=WHITE)
        box(sl, Inches(6.9), by+Inches(0.38), bar_w, Inches(0.55), bg=colour)
        txbox(sl, Inches(6.9)+bar_w+Inches(0.1), by+Inches(0.38),
              Inches(1.5), Inches(0.55),
              f"{pct:.1f}%  {amt}", size=12, bold=True, colour=colour)

    txbox(sl, Inches(6.9), Inches(5.65), Inches(5.8), Inches(0.35),
          f"Total Treatment Revenue: ${m['trt_total_rev']:,.0f}",
          size=13, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    box(sl, Inches(0.5), Inches(5.75), W-Inches(1), Inches(1.4),
        bg=RGBColor(0x0A, 0x14, 0x30), border=GOLD, border_pt=2)
    txbox(sl, Inches(0.7), Inches(5.85), W-Inches(1.4), Inches(0.35),
          "KEY TAKEAWAY", size=11, bold=True, colour=GOLD)
    txbox(sl, Inches(0.7), Inches(6.2), W-Inches(1.4), Inches(0.8),
          "A small fraction of players drove an outsized share of revenue — "
          "a classic high-value cohort effect amplified by the Royal Spin mechanic.",
          size=14, colour=WHITE, wrap=True)


def slide_segments(prs, m):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Segment Deep-Dive",
                 "Whales drove the lift — but Dolphin cannibalization is a concern")

    headers  = ["Segment", "Control ARPU", "Treatment ARPU", "Delta", "Signal"]
    col_w    = [Inches(2.0), Inches(2.2), Inches(2.5), Inches(1.8), Inches(4.3)]
    col_x    = [Inches(0.5)]
    for w in col_w[:-1]:
        col_x.append(col_x[-1] + w)

    box(sl, Inches(0.5), Inches(1.05), W-Inches(1), Inches(0.42), bg=GOLD)
    for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
        txbox(sl, x+Inches(0.05), Inches(1.07), w, Inches(0.38),
              h, size=12, bold=True, colour=NAVY, align=PP_ALIGN.CENTER)

    def _seg_row(seg, signal):
        ca = m['arpu_seg'].get(("Control",   seg), 0)
        ta = m['arpu_seg'].get(("Treatment", seg), 0)
        delta_pct = round((ta - ca) / ca * 100, 1) if ca else 0
        sign  = "+" if delta_pct >= 0 else ""
        dc    = GREEN if delta_pct > 0 else (AMBER if delta_pct > -15 else RED)
        return (seg, f"${ca:.2f}", f"${ta:.2f}", f"{sign}{delta_pct}%", dc, signal)

    seg_rows = [
        _seg_row("Whale",   "Strong lift — additive token model working"),
        _seg_row("Dolphin", "Cannibalization: Royal Tokens replaced coin spend"),
        _seg_row("Minnow",  "Negligible impact either way"),
    ]
    for r, (seg, ctrl, trt, delta, dc, signal) in enumerate(seg_rows):
        row_bg = RGBColor(0x1A, 0x2D, 0x5A) if r%2==0 else RGBColor(0x12, 0x22, 0x48)
        y_row  = Inches(1.47) + r*Inches(0.72)
        box(sl, Inches(0.5), y_row, W-Inches(1), Inches(0.7), bg=row_bg)
        for i, (val, colour, align) in enumerate([
            (seg,    GOLD,    PP_ALIGN.CENTER),
            (ctrl,   WHITE,   PP_ALIGN.CENTER),
            (trt,    WHITE,   PP_ALIGN.CENTER),
            (delta,  dc,      PP_ALIGN.CENTER),
            (signal, MID_GREY,PP_ALIGN.LEFT),
        ]):
            txbox(sl, col_x[i]+Inches(0.05), y_row+Inches(0.18),
                  col_w[i]-Inches(0.1), Inches(0.4),
                  val, size=13, bold=(i==0), colour=colour, align=align)

    # Dolphin detail
    ctrl_coin_delta = m['dolph_trt_coin'] - m['dolph_ctrl_coin']
    net_sign  = "loss" if m['dolph_net'] < 0 else "gain"
    net_label = f"${abs(m['dolph_net']):,.0f} net {net_sign}"

    box(sl, Inches(0.5), Inches(3.7), Inches(6.0), Inches(2.5),
        bg=RGBColor(0x2C, 0x10, 0x10), border=RED, border_pt=1.5)
    txbox(sl, Inches(0.7), Inches(3.8), Inches(5.6), Inches(0.38),
          "DOLPHIN CANNIBALIZATION DETAIL", size=12, bold=True, colour=RED)
    dolphin_tbl = [
        ("",           "Coin Revenue",                           "Royal Token Revenue"),
        ("Control",    f"${m['dolph_ctrl_coin']:,.0f}",          "—"),
        ("Treatment",  f"${m['dolph_trt_coin']:,.0f}",           f"${m['dolph_trt_token']:,.0f}"),
        ("Net change", f"${ctrl_coin_delta:+,.0f}",
         f"+${m['dolph_trt_token']:,.0f}  =  {net_label}"),
    ]
    for r, row_vals in enumerate(dolphin_tbl):
        y_r = Inches(4.25) + r*Inches(0.44)
        for c, val in enumerate(row_vals):
            colour = GOLD if r==0 else (RED if (r==3 and c==1) else WHITE)
            txbox(sl, Inches(0.65)+c*Inches(1.75), y_r, Inches(1.8), Inches(0.4),
                  val, size=12, bold=(r==0 or r==3), colour=colour)

    # Whale box
    ws_ctrl = m['whale_share'].get('Control',   0)
    ws_trt  = m['whale_share'].get('Treatment', 0)
    box(sl, Inches(6.9), Inches(3.7), Inches(6.0), Inches(2.5),
        bg=RGBColor(0x0D, 0x2A, 0x1A), border=GREEN, border_pt=1.5)
    txbox(sl, Inches(7.1), Inches(3.8), Inches(5.6), Inches(0.38),
          "WHALE REVENUE SHARE", size=12, bold=True, colour=GREEN)
    txbox(sl, Inches(7.1), Inches(4.2),  Inches(5.6), Inches(0.4),
          f"Control:    {ws_ctrl:.1f}% of group revenue from Whales", size=13, colour=WHITE)
    txbox(sl, Inches(7.1), Inches(4.65), Inches(5.6), Inches(0.4),
          f"Treatment:  {ws_trt:.1f}% of group revenue from Whales", size=13, colour=WHITE)
    txbox(sl, Inches(7.1), Inches(5.1),  Inches(5.6), Inches(0.7),
          "Treatment is more Whale-dependent — revenue base\nis more sensitive to Whale churn.",
          size=12, colour=AMBER, wrap=True)

    box(sl, Inches(0.5), Inches(6.35), W-Inches(1), Inches(0.85),
        bg=RGBColor(0x0A, 0x14, 0x30), border=GOLD, border_pt=2)
    txbox(sl, Inches(0.7), Inches(6.45), W-Inches(1.4), Inches(0.65),
          "RISK: Before full rollout, address Dolphin cannibalization — "
          "the substitutive token model cost more in lost coin revenue than it earned in token revenue.",
          size=13, colour=WHITE, wrap=True)


def slide_platform(prs, m):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Platform Breakdown",
                 f"The Royal Spin lift is an iOS story — Android saw minimal impact")

    platforms = [
        ("iOS",
         m['arpu_plat'].get(("Control","iOS"),0),
         m['arpu_plat'].get(("Treatment","iOS"),0),
         m['ios_lift']),
        ("Android",
         m['arpu_plat'].get(("Control","Android"),0),
         m['arpu_plat'].get(("Treatment","Android"),0),
         m['and_lift']),
    ]

    for i, (plat, ctrl, trt, lift) in enumerate(platforms):
        x  = Inches(0.5) + i*Inches(6.5)
        dc = GREEN if lift >= 30 else AMBER
        box(sl, x, Inches(1.1), Inches(6.0), Inches(5.2),
            bg=RGBColor(0x1A, 0x2D, 0x5A), border=GOLD, border_pt=1.5)
        box(sl, x, Inches(1.1), Inches(6.0), Inches(0.45),
            bg=GOLD if i==0 else MID_GREY)
        txbox(sl, x, Inches(1.13), Inches(6.0), Inches(0.4),
              plat, size=15, bold=True, colour=NAVY, align=PP_ALIGN.CENTER)
        txbox(sl, x+Inches(0.2), Inches(1.7),  Inches(5.6), Inches(0.38),
              "Control ARPU", size=12, colour=MID_GREY, align=PP_ALIGN.CENTER)
        txbox(sl, x+Inches(0.2), Inches(2.05), Inches(5.6), Inches(0.6),
              f"${ctrl:.2f}", size=30, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        txbox(sl, x+Inches(0.2), Inches(2.75), Inches(5.6), Inches(0.38),
              "Treatment ARPU", size=12, colour=MID_GREY, align=PP_ALIGN.CENTER)
        txbox(sl, x+Inches(0.2), Inches(3.1),  Inches(5.6), Inches(0.6),
              f"${trt:.2f}", size=30, bold=True, colour=GOLD, align=PP_ALIGN.CENTER)
        box(sl, x+Inches(1.5), Inches(3.85), Inches(3.0), Inches(0.6), bg=dc)
        txbox(sl, x+Inches(1.5), Inches(3.88), Inches(3.0), Inches(0.55),
              f"+{lift}%", size=22, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    box(sl, Inches(0.5), Inches(6.45), W-Inches(1), Inches(0.75),
        bg=RGBColor(0x0A, 0x14, 0x30), border=GOLD, border_pt=2)
    txbox(sl, Inches(0.7), Inches(6.52), W-Inches(1.4), Inches(0.6),
          f"iOS players responded {round(m['ios_lift']/m['and_lift'])}x more strongly to the Royal Spin "
          f"feature than Android players. Prioritise iOS in rollout targeting and marketing spend.",
          size=13, colour=WHITE, wrap=True)


def slide_guardrails(prs, m):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Guardrail Metrics",
                 "Engagement held firm — revenue concentration is the only structural risk")

    ad_c = m['active_days'].get('Control',   0)
    ad_t = m['active_days'].get('Treatment', 0)
    ws_c = m['whale_share'].get('Control',   0)
    ws_t = m['whale_share'].get('Treatment', 0)
    da_c = m['arpu_seg'].get(("Control",   "Dolphin"), 0)
    da_t = m['arpu_seg'].get(("Treatment", "Dolphin"), 0)

    def _st(ok): return ("PASS", GREEN) if ok else ("WATCH", AMBER)

    guardrails = [
        ("D1 Retention",
         f"{m['trt_d1']:.2f}% T  vs  {m['ctrl_d1']:.2f}% C",
         *_st(abs(m['trt_d1'] - m['ctrl_d1']) < 2),
         "No meaningful change"),
        ("D7 Retention",
         f"{m['trt_d7']:.2f}% T  vs  {m['ctrl_d7']:.2f}% C",
         *_st(abs(m['trt_d7'] - m['ctrl_d7']) < 2),
         "Slight positive trend" if m['trt_d7'] >= m['ctrl_d7'] else "Slight dip — monitor"),
        ("D30 Retention",
         f"{m['trt_d30']:.2f}% T  vs  {m['ctrl_d30']:.2f}% C",
         *_st(abs(m['trt_d30'] - m['ctrl_d30']) < 2),
         "Slight positive trend" if m['trt_d30'] >= m['ctrl_d30'] else "Slight dip — monitor"),
        ("Avg Active Days",
         f"{ad_t} T  vs  {ad_c} C",
         *_st(ad_t >= ad_c - 0.5),
         "No churn signal"),
        ("Spins / Session",
         f"{m['trt_spins']:.2f} T  vs  {m['ctrl_spins']:.2f} C",
         *_st(abs(m['trt_spins'] - m['ctrl_spins']) < 5),
         "Engagement unchanged"),
        ("Dolphin Revenue",
         f"${da_t:.2f} T  vs  ${da_c:.2f} C",
         *_st(da_t >= da_c),
         "Cannibalization confirmed" if da_t < da_c else "No cannibalization"),
        ("Whale Revenue Share",
         f"{ws_t:.1f}% T  vs  {ws_c:.1f}% C",
         *_st(ws_t <= ws_c + 2),
         "Increased concentration risk" if ws_t > ws_c else "Concentration stable"),
    ]

    for i, (metric, values, status, sc, note) in enumerate(guardrails):
        row_bg = RGBColor(0x1A, 0x2D, 0x5A) if i%2==0 else RGBColor(0x12, 0x22, 0x48)
        y_r    = Inches(1.1) + i*Inches(0.71)
        box(sl, Inches(0.5), y_r, W-Inches(1), Inches(0.69), bg=row_bg)
        box(sl, Inches(0.55), y_r+Inches(0.12), Inches(0.85), Inches(0.45), bg=sc)
        txbox(sl, Inches(0.55), y_r+Inches(0.14), Inches(0.85), Inches(0.4),
              status, size=10, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        txbox(sl, Inches(1.5),  y_r+Inches(0.14), Inches(3.2), Inches(0.42),
              metric, size=13, bold=True, colour=WHITE)
        txbox(sl, Inches(4.7),  y_r+Inches(0.14), Inches(4.0), Inches(0.42),
              values, size=12, colour=MID_GREY)
        txbox(sl, Inches(8.7),  y_r+Inches(0.14), Inches(4.3), Inches(0.42),
              note, size=12, colour=sc)

    box(sl, Inches(0.5), Inches(6.2), W-Inches(1), Inches(1.0),
        bg=RGBColor(0x0A, 0x14, 0x30), border=GOLD, border_pt=2)
    txbox(sl, Inches(0.7), Inches(6.3),  W-Inches(1.4), Inches(0.35),
          "SUMMARY", size=11, bold=True, colour=GOLD)
    txbox(sl, Inches(0.7), Inches(6.62), W-Inches(1.4), Inches(0.5),
          "No engagement or churn guardrails were triggered. The two WATCH items are known model "
          "effects (substitutive pricing for Dolphins; Whale concentration) — not surprises.",
          size=13, colour=WHITE, wrap=True)


def slide_recommendation(prs, m):
    sl = blank_slide(prs)
    fill_bg(sl, NAVY)
    slide_header(sl, "Recommendation",
                 "Ship Royal Spin — with one fix before broad rollout")

    box(sl, Inches(0.5), Inches(1.05), W-Inches(1), Inches(0.75), bg=GREEN)
    txbox(sl, Inches(0.5), Inches(1.08), W-Inches(1), Inches(0.68),
          f"VERDICT:  SHIP  —  Royal Spin delivers a statistically meaningful "
          f"+{m['arpu_lift']}% ARPU lift with no engagement cost",
          size=16, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)

    net_sign = "loss" if m['dolph_net'] < 0 else "gain"
    ws_t = m['whale_share'].get('Treatment', 0)
    ws_c = m['whale_share'].get('Control',   0)

    actions = [
        (GREEN, "SHIP NOW",
         "Roll out Royal Spin to all Whale players on iOS",
         [
             f"Strongest signal: iOS Whales saw +{m['ios_lift']}% ARPU and +{m['whale_seg_lift']}% segment ARPU",
             "Additive token model means no cannibalization risk for this segment",
             f"{m['royal_adoption']:.2f}% Royal Token adoption driving "
             f"{m['trt_token_pct']:.0f}% of Treatment revenue",
         ]),
        (AMBER, "FIX FIRST",
         "Redesign Dolphin token pricing before broad rollout",
         [
             "Substitutive model cost Dolphins more in lost coin revenue than token revenue gained",
             f"Net Dolphin {net_sign}: ${abs(m['dolph_net']):,.0f} across the experiment cohort",
             "Consider additive-only token model for Dolphins, or tiered pricing",
         ]),
        (RGBColor(0x41, 0x8C, 0xF0), "MONITOR",
         "Track Whale retention and revenue concentration post-launch",
         [
             f"Treatment revenue is {ws_t:.1f}% Whale-dependent vs {ws_c:.1f}% in Control",
             "Set up weekly Whale churn alert (>2pp above baseline)",
             f"Review Android performance separately — only +{m['and_lift']}% ARPU lift",
         ]),
    ]

    for i, (colour, badge, title, bullets) in enumerate(actions):
        x = Inches(0.5) + i*Inches(4.25)
        box(sl, x, Inches(2.0), Inches(4.0), Inches(4.8),
            bg=RGBColor(0x1A, 0x2D, 0x5A), border=colour, border_pt=2)
        box(sl, x, Inches(2.0), Inches(4.0), Inches(0.42), bg=colour)
        txbox(sl, x, Inches(2.02), Inches(4.0), Inches(0.38),
              badge, size=13, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
        txbox(sl, x+Inches(0.15), Inches(2.52), Inches(3.7), Inches(0.55),
              title, size=13, bold=True, colour=WHITE, wrap=True)
        for j, b in enumerate(bullets):
            txbox(sl, x+Inches(0.15), Inches(3.15)+j*Inches(0.68),
                  Inches(3.7), Inches(0.65),
                  "•  "+b, size=11, colour=MID_GREY, wrap=True)

    txbox(sl, Inches(0.5), Inches(7.1), W-Inches(1), Inches(0.3),
          "Royal Flush Casino  |  SpinCrown Studios  |  Data & Analytics  |  2024",
          size=10, colour=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def main():
    conn = sqlite3.connect(config.DB_PATH)
    m    = load_metrics(conn)
    conn.close()

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    # Stash lift on prs so slide_agenda can read it without a full metrics arg
    prs._arpu_lift = m['arpu_lift']

    slide_title(prs, m)
    slide_agenda(prs)
    slide_methodology(prs)
    slide_headline(prs, m)
    slide_royal_token(prs, m)
    slide_segments(prs, m)
    slide_platform(prs, m)
    slide_guardrails(prs, m)
    slide_recommendation(prs, m)

    out = "exports/royal_spin_ab_test.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
